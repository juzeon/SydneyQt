import json
import re

import pptx
from pptx.enum.shapes import MSO_SHAPE_TYPE


def read_pptx_text(filename):
    def check_recursively_for_text(this_set_of_shapes, text_run):
        for shape in this_set_of_shapes:
            if shape.shape_type == MSO_SHAPE_TYPE.GROUP:
                check_recursively_for_text(shape.shapes, text_run)
            else:
                if hasattr(shape, "text"):
                    text_run.append(shape.text)
        return text_run

    output = {}
    presentation = pptx.Presentation(filename)
    for i, slide in enumerate(presentation.slides):
        slide_number = i + 1
        text_run = []
        text_run = check_recursively_for_text(slide.shapes, text_run)
        output[f'Slide {slide_number}'] = re.sub('\\s+', ' ', ' '.join(text_run).strip())
    return json.dumps(output, ensure_ascii=False)
